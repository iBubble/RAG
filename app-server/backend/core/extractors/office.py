import logging
from typing import List
from pathlib import Path

logger = logging.getLogger(__name__)

from .utils import _table_to_markdown

def _extract_unstructured(file_path: str) -> str:
    """
    使用 Unstructured 库进行精细结构化提取。
    WHY: 对于 PPTX/DOCX/XLSX 等复杂排版文档，Unstructured 能识别标题层级、
         表格边界和列表结构，产出的 Knowledge Chunk 天然附带语义上下文，
         大幅提升 RAG 检索问答的精准度。
    适用: PPTX（55份，主力）、DOCX（25份）、XLSX（8份）。
    """
    from pathlib import Path

    path = Path(file_path)
    suffix = path.suffix.lower()
    logger.info(f"Unstructured 结构化解析: {path.name}")

    # WHY: 针对大型工程 Excel 表，直接绕过 Unstructured（其会提取全量导致超时）。
    #      强制走原生的 _extract_xlsx/_extract_xls 获取行数大纲，完整数据交由表格库处理。
    if suffix in (".xlsx", ".xls"):
        logger.info(f"{path.name} 为大型表格格式，绕过 Unstructured 强制走精简解析")
        return _fallback_extract(file_path, suffix)

    try:
        from unstructured.partition.auto import partition

        elements = partition(filename=str(path))

        if not elements:
            logger.warning(f"Unstructured 未提取到内容: {path.name}，回退原生解析")
            return _fallback_extract(file_path, suffix)

        # 将 Unstructured 的元素转为结构化纯文本
        texts = []
        for el in elements:
            el_type = type(el).__name__
            text = str(el).strip()
            if not text:
                continue

            # 为标题和表格添加语义标记，帮助大模型理解文档结构
            if el_type == "Title":
                texts.append(f"## {text}")
            elif el_type == "Table":
                texts.append(f"[表格]\n{text}")
            elif el_type == "ListItem":
                texts.append(f"- {text}")
            else:
                texts.append(text)

        result = "\n\n".join(texts)
        logger.info(f"Unstructured 提取完成: {path.name}，{len(result)} 字符，{len(elements)} 个元素")
        return result

    except Exception as e:
        logger.error(f"Unstructured 解析失败 ({path.name}): {e}，回退原生解析")
        return _fallback_extract(file_path, suffix)


def _fallback_extract(file_path: str, suffix: str) -> str:
    """
    当 Unstructured 引擎出错时，回退到原生的 python-pptx / python-docx 等基础提取器。
    WHY: 绝不能因为一个库的异常就让整个解析链路断裂。
    """
    fallback_map = {
        ".pptx": _extract_pptx,
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
    }
    fallback = fallback_map.get(suffix)
    if fallback:
        logger.info(f"回退到原生解析器: {suffix}")
        return fallback(file_path)
    return ""


def _extract_doc(file_path: str) -> str:
    """
    旧版 .doc 格式提取。
    优先尝试 antiword（Linux 常用轻量提取器），如果在 macOS 失败尝试 textutil。
    """
    import subprocess
    
    # 尝试 antiword
    try:
        result = subprocess.run(
            ["antiword", file_path],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return result.stdout
        else:
            logger.warning(f"antiword 提取 .doc 失败: {result.stderr}")
    except FileNotFoundError:
        pass

    # 尝试 textutil (macOS fallback)
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-output", tmp_path, file_path],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return Path(tmp_path).read_text(encoding="utf-8", errors="ignore")
    except FileNotFoundError:
        pass
    finally:
        Path(tmp_path).unlink(missing_ok=True)
        
    logger.error("系统中未找到 antiword 且 textutil 不可用，.doc 解析失败")
    return ""


def _extract_docx(file_path: str) -> str:
    """使用 python-docx 提取 Word 文档文本（含表格）。"""
    from docx import Document

    doc = Document(file_path)
    texts = []

    for para in doc.paragraphs:
        if para.text.strip():
            texts.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                texts.append(row_text)

    # 提取嵌入图片
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                img_bytes = rel.target_part.blob
                import tempfile
                from pathlib import Path
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    tmp.write(img_bytes)
                    tmp_path = tmp.name
                
                from .vision_extractor import extract_image_vision
                vision_text = extract_image_vision(tmp_path)
                
                Path(tmp_path).unlink(missing_ok=True)
                if vision_text and len(vision_text) > 10:
                    texts.append(f"[嵌入图片内容]\n{vision_text}")
            except Exception as e:
                logger.debug(f"DOCX 提取图片异常: {e}")

    return "\n\n".join(texts)


def _smart_extract_sheet(raw_matrix: list, max_row: int, max_col: int,
                         sheet_name: str = "") -> str:
    """
    通用的智能工作表提取器：处理合并单元格展开后的原始矩阵。
    WHY: 工程类统计报表（国土调查/规划编制）普遍使用多行合并表头，
         简单逐行拼接会导致列位移，大模型无法将数值与正确的列名关联。
    策略:
      1. 自动检测标题行、单位行、表头行、数据行
      2. 多行表头合并为单行（仅保留最后 2 级，避免冗长）
      3. 数据行采用"列名: 值"键值对格式，省略零值列（方案 B）
    """
    # ── 第 1 步：扫描识别行类型 ──
    header_keywords = {"行政区域", "名称", "代码", "面积", "编号", "序号",
                       "类型", "地类", "用途", "权属", "坐落", "单位名称",
                       "项目", "单位", "任务量", "投资", "工程量", "金额", "合计", "小计", "总计"}
    title_row = None
    unit_row = None
    header_start = None
    header_end = None

    scan_limit = min(20, max_row)
    for idx in range(scan_limit):
        row_str = " ".join(str(v) for v in raw_matrix[idx] if v is not None)
        if not row_str.strip():
            continue
        if any(kw in row_str for kw in ("汇总表", "统计表", "一览表", "明细表", "情况表")):
            title_row = idx
        elif "单位" in row_str and any(u in row_str for u in ("公顷", "亩", "平方", "万元", "元")) and len([str(v).strip() for v in raw_matrix[idx] if v is not None and str(v).strip()]) <= 3:
            unit_row = idx
        elif any(kw in row_str for kw in header_keywords):
            if header_start is None:
                header_start = idx
            header_end = idx
        elif header_start is not None and header_end is not None:
            # 表头检测结束，遇到非表头行即停止
            break

    # 如果未检测到复合表头，回退为简单逐行提取
    if header_start is None:
        return _simple_row_extract(raw_matrix, max_row, max_col, sheet_name)

    data_start = header_end + 1

    # ── 第 2 步：合并多行表头为单行 ──
    merged_headers = []
    for col_idx in range(max_col):
        parts = []
        for row_idx in range(header_start, header_end + 1):
            val = raw_matrix[row_idx][col_idx]
            if val is not None:
                s = str(val).replace("\n", "").strip()
                if s and s not in parts:
                    parts.append(s)
        # WHY: 只保留最后 2 层，避免"土地利用现状分类面积汇总表_单位：公顷_行政区域_名称"这样的冗长路径
        if len(parts) > 2:
            parts = parts[-2:]
        merged_headers.append("_".join(parts) if parts else "")

    # 找出有效列（表头非空的列）
    valid_cols = [i for i, h in enumerate(merged_headers) if h]

    # ── 第 3 步：组装输出文本 ──
    texts = []

    # 标题与单位
    if sheet_name:
        texts.append(f"[工作表: {sheet_name}]")
    if title_row is not None:
        title_val = next((str(v).strip() for v in raw_matrix[title_row] if v is not None and str(v).strip()), "")
        if title_val:
            texts.append(f"[{title_val}]")
    if unit_row is not None:
        unit_val = next((str(v).strip() for v in raw_matrix[unit_row] if v is not None and str(v).strip()), "")
        if unit_val:
            texts.append(f"[{unit_val}]")

    # 表头总览（供大模型理解整体结构）
    texts.append("字段结构: " + " | ".join(merged_headers[i] for i in valid_cols))

    # 数据行精简版：最多展开所有有效行（上限 10000），将键值对拼接到文本中
    row_count = max_row - data_start
    if row_count > 0:
        texts.append(f"数据行数: 共 {row_count} 行")
        
        limit = max_row
        for row_idx in range(data_start, limit):
            row_items = []
            for col_idx in valid_cols:
                val = raw_matrix[row_idx][col_idx]
                if val is not None:
                    s = str(val).strip()
                    if s and s.lower() not in ("nan", "none"):
                        h = merged_headers[col_idx]
                        if h:
                            row_items.append(f"{h}: {s}")
                        else:
                            row_items.append(s)
            if row_items:
                texts.append(" | ".join(row_items))

    return "\n".join(texts)


def _simple_row_extract(raw_matrix: list, max_row: int, max_col: int,
                        sheet_name: str = "") -> str:
    """
    简单逐行提取（无复合表头的普通表格回退方案）。
    WHY: 并非所有 Excel 都有复合表头，简单表格无需键值对化，保持轻量输出。
    """
    texts = []
    if sheet_name:
        texts.append(f"[工作表: {sheet_name}]")

    # 尝试找到第一行非空的作为表头结构
    headers = []
    header_row_idx = -1
    for row_idx in range(min(5, max_row)):
        row_values = []
        for col_idx in range(max_col):
            val = raw_matrix[row_idx][col_idx]
            if val is not None:
                s = str(val).strip()
                if s and s.lower() not in ("nan", "none"):
                    row_values.append(s)
        if row_values:
            headers = row_values
            header_row_idx = row_idx
            break
            
    if headers:
        texts.append("字段结构: " + " | ".join(headers))
        
    texts.append(f"数据行数: 共 {max_row} 行")

    data_start = header_row_idx + 1 if header_row_idx >= 0 else 0
    limit = max_row
    for row_idx in range(data_start, limit):
        row_values = []
        for col_idx in range(max_col):
            val = raw_matrix[row_idx][col_idx]
            if val is not None:
                s = str(val).strip()
                if s and s.lower() not in ("nan", "none"):
                    if headers and col_idx < len(headers) and headers[col_idx]:
                        row_values.append(f"{headers[col_idx]}: {s}")
                    else:
                        row_values.append(s)
        if row_values:
            texts.append(" | ".join(row_values))

    return "\n".join(texts)


def _parse_merged_cells_from_zip(file_path: str) -> dict:
    """
    从 .xlsx ZIP 包的 XML 中直接解析每个 Sheet 的合并单元格范围。
    WHY: openpyxl 的 read_only=True 模式不暴露 merged_cells 属性，
         但 read_only 模式加载速度快 100 倍（0.06s vs 数分钟）。
         合并信息存储在每个 sheet 的 XML 中，直接解析零开销。
    返回: {sheet_name: [(min_row, min_col, max_row, max_col), ...]}
    """
    import zipfile
    from xml.etree import ElementTree as ET

    result = {}
    ns = '{http://schemas.openxmlformats.org/spreadsheetml/2006/main}'

    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            # 解析 workbook.xml 获取 sheet 名称顺序
            wb_xml = ET.parse(zf.open('xl/workbook.xml'))
            sheets_el = wb_xml.findall(f'{ns}sheets/{ns}sheet')
            sheet_names = [s.get('name') for s in sheets_el]

            for idx, name in enumerate(sheet_names):
                sheet_path = f'xl/worksheets/sheet{idx + 1}.xml'
                if sheet_path not in zf.namelist():
                    continue

                sheet_xml = ET.parse(zf.open(sheet_path))
                merge_cells = sheet_xml.findall(f'{ns}mergeCells/{ns}mergeCell')

                ranges = []
                for mc in merge_cells:
                    ref = mc.get('ref')  # 例如 "A1:C3"
                    if not ref or ':' not in ref:
                        continue
                    from openpyxl.utils import range_boundaries
                    min_c, min_r, max_c, max_r = range_boundaries(ref)
                    ranges.append((min_r, min_c, max_r, max_c))

                if ranges:
                    result[name] = ranges
    except Exception as e:
        logger.warning(f"解析合并单元格 XML 失败（将跳过合并展开）: {e}")

    return result


def _build_safe_matrix(ws, sheet_name: str):
    """
    从 worksheet 流式读取数据并构建安全修剪后的矩阵。
    WHY: 抽取公共逻辑，避免 _extract_xlsx 和 _extract_tables_xlsx 重复代码。
    返回: (raw_matrix, actual_max_row, actual_max_col) 或 None
    """
    MAX_SAFE_ROWS = 10000
    MAX_SAFE_COLS = 200

    raw_data = []
    for idx, row in enumerate(ws.iter_rows(values_only=True)):
        if idx >= MAX_SAFE_ROWS:
            logger.warning(
                f"[防 OOM 熔断] {sheet_name} 行数超过 {MAX_SAFE_ROWS}，执行安全截断"
            )
            break
        raw_data.append(row[:MAX_SAFE_COLS])

    if not raw_data:
        return None

    # 动态修剪真实边界（去掉全 None 的尾部行/列）
    actual_max_row = 0
    actual_max_col = 0

    for r_idx, row in enumerate(raw_data):
        last_valid_c = -1
        for c_idx in range(len(row) - 1, -1, -1):
            if row[c_idx] is not None and str(row[c_idx]).strip() != "":
                last_valid_c = c_idx
                break

        if last_valid_c >= 0:
            actual_max_row = r_idx + 1
            if last_valid_c + 1 > actual_max_col:
                actual_max_col = last_valid_c + 1

    if actual_max_row == 0 or actual_max_col == 0:
        return None

    # 构建真实大小的矩阵
    raw_matrix = []
    for r_idx in range(actual_max_row):
        row = raw_data[r_idx]
        clamped_row = list(row[:actual_max_col])
        if len(clamped_row) < actual_max_col:
            clamped_row.extend([None] * (actual_max_col - len(clamped_row)))
        raw_matrix.append(clamped_row)

    return raw_matrix, actual_max_row, actual_max_col


def _apply_merged_cells(raw_matrix, merge_ranges, actual_max_row, actual_max_col):
    """
    将合并单元格的值填充到整个合并区域。
    WHY: 抽取公共逻辑，合并展开在矩阵构建之后统一执行。
    """
    for min_r, min_c, max_r, max_c in merge_ranges:
        max_r = min(max_r, actual_max_row)
        max_c = min(max_c, actual_max_col)

        if min_r > actual_max_row or min_c > actual_max_col:
            continue

        val = raw_matrix[min_r - 1][min_c - 1]
        for r in range(min_r, max_r + 1):
            for c in range(min_c, max_c + 1):
                raw_matrix[r - 1][c - 1] = val


def _extract_xlsx(file_path: str) -> str:
    """
    使用 openpyxl 提取 .xlsx 表格数据（支持复合表头智能合并）。
    WHY: 使用 read_only=True 极速加载（0.06s），合并信息从 ZIP XML 解析。
         彻底解决大 Excel 文件卡死和 OOM 问题。
    """
    from openpyxl import load_workbook

    # WHY: 先从 ZIP XML 解析合并单元格（毫秒级），再用 read_only 读数据
    merged_map = _parse_merged_cells_from_zip(file_path)

    wb = load_workbook(file_path, read_only=True, data_only=True)
    all_texts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]

        result = _build_safe_matrix(ws, sheet_name)
        if result is None:
            continue
        raw_matrix, actual_max_row, actual_max_col = result

        # 展开合并单元格
        if sheet_name in merged_map:
            _apply_merged_cells(
                raw_matrix, merged_map[sheet_name],
                actual_max_row, actual_max_col,
            )

        text = _smart_extract_sheet(
            raw_matrix, actual_max_row, actual_max_col, sheet_name,
        )
        if text.strip():
            all_texts.append(text)

    wb.close()
    return "\n\n".join(all_texts)


def _extract_xls(file_path: str) -> str:
    """
    使用 xlrd 提取旧版 .xls 表格数据（支持复合表头智能合并）。
    WHY: 旧版 .xls 同样存在合并单元格问题，需与 .xlsx 保持一致的提取质量。
    注意: xlrd 不直接提供合并单元格 API，使用 merged_cells 属性获取合并区域。
    """
    import xlrd

    wb = xlrd.open_workbook(file_path, formatting_info=True)
    all_texts = []

    for sheet in wb.sheets():
        max_row = sheet.nrows
        max_col = sheet.ncols
        if not max_row or not max_col:
            continue

        # 构建原始矩阵
        raw_matrix = [[None] * max_col for _ in range(max_row)]
        for rx in range(max_row):
            for cx in range(max_col):
                val = sheet.cell_value(rx, cx)
                if val != "":
                    raw_matrix[rx][cx] = val

        # 展开合并单元格
        for (rlo, rhi, clo, chi) in sheet.merged_cells:
            val = sheet.cell_value(rlo, clo)
            for r in range(rlo, rhi):
                for c in range(clo, chi):
                    raw_matrix[r][c] = val

        text = _smart_extract_sheet(raw_matrix, max_row, max_col, sheet.name)
        if text.strip():
            all_texts.append(text)

    return "\n\n".join(all_texts)


def _extract_pptx(file_path: str) -> str:
    """使用 python-pptx 提取 PPT 幻灯片文本。"""
    from pptx import Presentation

    prs = Presentation(file_path)
    texts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[幻灯片 {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text.strip())

            # 提取表格
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells
                        if cell.text.strip()
                    )
                    if row_text:
                        slide_texts.append(row_text)

        if len(slide_texts) > 1:
            texts.append("\n".join(slide_texts))

    return "\n\n".join(texts)


def _extract_tables_pptx(file_path: str) -> List[dict]:
    """
    从 .pptx 文件提取所有结构化表格。
    WHY: PPT 演示文稿中嵌入的数据表格（如项目概况表、投资估算表）
         同样需要完整注册到表格注册表，供报告生成时精确注入。
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    tables = []
    filename = Path(file_path).name

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_table:
                continue

            tbl = shape.table
            if len(tbl.rows) < 2:
                continue

            # 第一行作为表头
            headers = []
            for cell in tbl.rows[0].cells:
                h = cell.text.strip().replace("\n", " ")
                # WHY: 相邻合并单元格会产生重复列名，去重
                if not headers or h != headers[-1]:
                    headers.append(h)

            if not any(h for h in headers):
                continue

            # 提取数据行
            data_rows = []
            for row in list(tbl.rows)[1:]:
                cells = []
                for cell in row.cells:
                    c = cell.text.strip().replace("\n", " ")
                    if not cells or c != cells[-1]:
                        cells.append(c)
                # 补齐到 headers 长度
                while len(cells) < len(headers):
                    cells.append("")
                if any(c for c in cells):
                    data_rows.append(cells[:len(headers)])

            if len(data_rows) < 1:
                continue

            # 尝试从幻灯片标题推断表格名称
            title = f"幻灯片{slide_idx}_表格"
            for s in slide.shapes:
                if s.has_text_frame and s.text.strip():
                    candidate = s.text.strip()
                    if len(candidate) < 60:
                        title = candidate
                        break

            md = _table_to_markdown(title, "", headers, data_rows)

            tables.append({
                "title": title,
                "sheet_name": f"幻灯片{slide_idx}",
                "headers": headers,
                "rows": data_rows,
                "unit": "",
                "markdown": md,
                "source_file": filename,
                "row_count": len(data_rows),
                "char_count": len(md),
            })

    logger.info(f"[表格提取] {filename}: 提取到 {len(tables)} 张 PPT 表格")
    return tables

def _extract_tables_xlsx(file_path: str) -> List[dict]:
    """
    从 .xlsx 文件提取所有结构化表格。
    WHY: 使用 read_only=True 极速加载 + ZIP XML 解析合并信息，防 OOM。
    """
    from openpyxl import load_workbook

    merged_map = _parse_merged_cells_from_zip(file_path)

    wb = load_workbook(file_path, read_only=True, data_only=True)
    tables = []
    filename = Path(file_path).name

    header_keywords = {"行政区域", "名称", "代码", "面积", "编号", "序号",
                       "类型", "地类", "用途", "权属", "坐落", "单位名称",
                       "乡镇", "项目", "产量", "产值", "数量"}

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]

        result = _build_safe_matrix(ws, sheet_name)
        if result is None:
            continue
        raw_matrix, actual_max_row, actual_max_col = result

        if actual_max_row < 3:
            continue

        # 展开合并单元格
        if sheet_name in merged_map:
            _apply_merged_cells(
                raw_matrix, merged_map[sheet_name],
                actual_max_row, actual_max_col,
            )

        # 检测表头和标题
        title = ""
        unit = ""
        header_start = None
        header_end = None

        scan_limit = min(20, actual_max_row)
        for idx in range(scan_limit):
            row_str = " ".join(str(v) for v in raw_matrix[idx] if v is not None)
            if not row_str.strip():
                continue
            if any(kw in row_str for kw in ("汇总表", "统计表", "一览表", "明细表",
                                            "计算表", "成果表", "调查表", "情况表")):
                title = next(
                    (str(v).strip() for v in raw_matrix[idx] if v is not None and str(v).strip()),
                    sheet_name
                )
            elif "单位" in row_str and any(u in row_str for u in ("公顷", "亩", "平方", "万元", "元", "吨", "kg")) and len([str(v).strip() for v in raw_matrix[idx] if v is not None and str(v).strip()]) <= 3:
                unit = next(
                    (str(v).strip() for v in raw_matrix[idx] if v is not None and str(v).strip()),
                    ""
                )
            elif any(kw in row_str for kw in header_keywords):
                # WHY: 防止合并标题行被误识别为表头 —
                #      当所有非空列值相同时（如"泸县高标准农田..."×5），
                #      这是合并标题行而非数据表头，跳过。
                non_none_vals = [str(v).strip() for v in raw_matrix[idx] if v is not None and str(v).strip()]
                if non_none_vals and len(set(non_none_vals)) == 1:
                    continue  # 合并标题行，不是表头
                if header_start is None:
                    header_start = idx
                header_end = idx
            elif header_start is not None and header_end is not None:
                break

        if header_start is None:
            # 回退：用第一行作为表头
            first_non_empty = 0
            for idx in range(min(5, actual_max_row)):
                row_str = " ".join(str(v) for v in raw_matrix[idx] if v is not None)
                if row_str.strip():
                    first_non_empty = idx
                    break
            header_start = first_non_empty
            header_end = first_non_empty

        # 合并多行表头
        merged_headers = []
        for col_idx in range(actual_max_col):
            parts = []
            for row_idx in range(header_start, header_end + 1):
                val = raw_matrix[row_idx][col_idx]
                if val is not None:
                    s = str(val).replace("\n", "").strip()
                    if s and s not in parts:
                        parts.append(s)
            if len(parts) > 2:
                parts = parts[-2:]
            merged_headers.append("_".join(parts) if parts else "")

        valid_cols = [i for i, h in enumerate(merged_headers) if h]
        if not valid_cols:
            continue

        headers = [merged_headers[i] for i in valid_cols]
        data_start = header_end + 1

        # 提取数据行
        data_rows = []
        for row_idx in range(data_start, actual_max_row):
            cells = []
            has_data = False
            for col_idx in valid_cols:
                val = raw_matrix[row_idx][col_idx]
                if val is None:
                    cells.append("")
                else:
                    s = str(val).strip()
                    if s.lower() in ("nan", "none"):
                        cells.append("")
                    else:
                        cells.append(s)
                        if s:
                            has_data = True
            if has_data:
                data_rows.append(cells)

        # 过滤掉行数太少的"假表格"
        if len(data_rows) < 2:
            continue

        # 找出具有跨行合并的列名
        merged_cols = []
        if sheet_name in merged_map:
            merged_col_indices = set()
            for min_r, min_c, max_r, max_c in merged_map[sheet_name]:
                if max_r > min_r:
                    for col_idx in range(min_c, max_c + 1):
                        merged_col_indices.add(col_idx - 1)
            for c_idx in valid_cols:
                if c_idx in merged_col_indices:
                    h_name = merged_headers[c_idx]
                    if h_name and h_name not in merged_cols:
                        merged_cols.append(h_name)

        if not title:
            title = sheet_name

        md = _table_to_markdown(title, unit, headers, data_rows)

        tables.append({
            "title": title,
            "sheet_name": sheet_name,
            "headers": headers,
            "rows": data_rows,
            "unit": unit,
            "markdown": md,
            "source_file": filename,
            "row_count": len(data_rows),
            "char_count": len(md),
            "merged_cols": merged_cols,
        })

    wb.close()
    logger.info(f"[表格提取] {filename}: 提取到 {len(tables)} 张表格")
    return tables


def _extract_tables_docx(file_path: str) -> List[dict]:
    """
    从 .docx 文件提取所有结构化表格。
    WHY: Word 报告中嵌入的数据表格同样需要完整注册。
    """
    from docx import Document
    import re

    doc = Document(file_path)
    tables = []
    filename = Path(file_path).name

    # 预提取段落文本，用于表格标题匹配
    paragraphs = [p.text.strip() for p in doc.paragraphs]

    for tbl_idx, table in enumerate(doc.tables):
        if len(table.rows) < 3:
            continue

        # 提取表头（第一行）
        headers = [cell.text.strip().replace("\n", " ") for cell in table.rows[0].cells]
        # 去重相邻合并单元格产生的重复列名
        deduped = []
        for h in headers:
            if not deduped or h != deduped[-1]:
                deduped.append(h)
        headers = deduped if deduped else headers

        if not any(h for h in headers):
            continue

        # 提取数据行（跳过第一行表头）
        data_rows = []
        for row in table.rows[1:]:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            # 去重相邻合并单元格
            deduped_cells = []
            for c in cells:
                if not deduped_cells or c != deduped_cells[-1]:
                    deduped_cells.append(c)
            # 补齐到 headers 长度
            while len(deduped_cells) < len(headers):
                deduped_cells.append("")
            if any(c for c in deduped_cells):
                data_rows.append(deduped_cells[:len(headers)])

        if len(data_rows) < 2:
            continue

        # 尝试从前方段落找表格标题
        title = f"表{tbl_idx + 1}"
        table_title_pattern = re.compile(r'表\s*\d+|表格\s*\d+')
        # 用 table 的 XML 元素在文档中定位
        tbl_element = table._tbl
        parent = tbl_element.getparent()
        if parent is not None:
            prev = tbl_element.getprevious()
            if prev is not None and prev.text:
                candidate = prev.text.strip()
                if candidate and len(candidate) < 80:
                    title = candidate

        unit = ""
        md = _table_to_markdown(title, unit, headers, data_rows)

        tables.append({
            "title": title,
            "sheet_name": "",
            "headers": headers,
            "rows": data_rows,
            "unit": unit,
            "markdown": md,
            "source_file": filename,
            "row_count": len(data_rows),
            "char_count": len(md),
        })

    logger.info(f"[表格提取] {filename}: 提取到 {len(tables)} 张表格")
    return tables


def _extract_tables_xls(file_path: str) -> List[dict]:
    """
    从 .xls 文件提取结构化表格（与 xlsx 逻辑对齐）。
    WHY: 旧版 .xls 格式在工程行业仍大量使用。
    """
    import xlrd

    wb = xlrd.open_workbook(file_path, formatting_info=True)
    tables = []
    filename = Path(file_path).name

    header_keywords = {"行政区域", "名称", "代码", "面积", "编号", "序号",
                       "类型", "地类", "用途", "权属", "坐落", "单位名称",
                       "乡镇", "项目", "产量", "产值", "数量"}

    for sheet in wb.sheets():
        max_row = sheet.nrows
        max_col = sheet.ncols
        if not max_row or not max_col or max_row < 3:
            continue

        raw_matrix = [[None] * max_col for _ in range(max_row)]
        for rx in range(max_row):
            for cx in range(max_col):
                val = sheet.cell_value(rx, cx)
                if val != "":
                    raw_matrix[rx][cx] = val

        for crange in sheet.merged_cells:
            rlo, rhi, clo, chi = crange
            val = raw_matrix[rlo][clo]
            for r in range(rlo, rhi):
                for c in range(clo, chi):
                    raw_matrix[r][c] = val

        # 复用 xlsx 相同的检测逻辑
        title = ""
        unit = ""
        header_start = None
        header_end = None

        scan_limit = min(20, max_row)
        for idx in range(scan_limit):
            row_str = " ".join(str(v) for v in raw_matrix[idx] if v is not None)
            if not row_str.strip():
                continue
            if any(kw in row_str for kw in ("汇总表", "统计表", "一览表", "明细表",
                                            "计算表", "成果表", "调查表", "情况表")):
                title = next(
                    (str(v).strip() for v in raw_matrix[idx] if v is not None and str(v).strip()),
                    sheet.name
                )
            elif "单位" in row_str and any(u in row_str for u in ("公顷", "亩", "平方", "万元", "元", "吨")) and len([str(v).strip() for v in raw_matrix[idx] if v is not None and str(v).strip()]) <= 3:
                unit = next(
                    (str(v).strip() for v in raw_matrix[idx] if v is not None and str(v).strip()),
                    ""
                )
            elif any(kw in row_str for kw in header_keywords):
                # WHY: 防止合并标题行被误识别为表头（同 _extract_tables_xlsx）
                non_none_vals = [str(v).strip() for v in raw_matrix[idx] if v is not None and str(v).strip()]
                if non_none_vals and len(set(non_none_vals)) == 1:
                    continue
                if header_start is None:
                    header_start = idx
                header_end = idx
            elif header_start is not None and header_end is not None:
                break

        if header_start is None:
            first_non_empty = 0
            for idx in range(min(5, max_row)):
                row_str = " ".join(str(v) for v in raw_matrix[idx] if v is not None)
                if row_str.strip():
                    first_non_empty = idx
                    break
            header_start = first_non_empty
            header_end = first_non_empty

        merged_headers = []
        for col_idx in range(max_col):
            parts = []
            for row_idx in range(header_start, header_end + 1):
                val = raw_matrix[row_idx][col_idx]
                if val is not None:
                    s = str(val).replace("\n", "").strip()
                    if s and s not in parts:
                        parts.append(s)
            if len(parts) > 2:
                parts = parts[-2:]
            merged_headers.append("_".join(parts) if parts else "")

        valid_cols = [i for i, h in enumerate(merged_headers) if h]
        if not valid_cols:
            continue

        headers = [merged_headers[i] for i in valid_cols]
        data_start = header_end + 1

        data_rows = []
        for row_idx in range(data_start, max_row):
            cells = []
            has_data = False
            for col_idx in valid_cols:
                val = raw_matrix[row_idx][col_idx]
                if val is None:
                    cells.append("")
                else:
                    s = str(val).strip()
                    cells.append(s)
                    if s:
                        has_data = True
            if has_data:
                data_rows.append(cells)

        if len(data_rows) < 2:
            continue

        # 找出具有跨行合并的列名
        merged_cols = []
        if sheet.merged_cells:
            merged_col_indices = set()
            for rlo, rhi, clo, chi in sheet.merged_cells:
                if rhi > rlo + 1:
                    for col_idx in range(clo, chi):
                        merged_col_indices.add(col_idx)
            for c_idx in valid_cols:
                if c_idx in merged_col_indices:
                    h_name = merged_headers[c_idx]
                    if h_name and h_name not in merged_cols:
                        merged_cols.append(h_name)

        if not title:
            title = sheet.name

        md = _table_to_markdown(title, unit, headers, data_rows)
        tables.append({
            "title": title,
            "sheet_name": sheet.name,
            "headers": headers,
            "rows": data_rows,
            "unit": unit,
            "markdown": md,
            "source_file": filename,
            "row_count": len(data_rows),
            "char_count": len(md),
            "merged_cols": merged_cols,
        })

    logger.info(f"[表格提取] {filename}: 提取到 {len(tables)} 张表格")
    return tables


